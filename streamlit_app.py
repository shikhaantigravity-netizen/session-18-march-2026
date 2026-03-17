import streamlit as st
import os
import json
import random
import sqlite3
import pandas as pd
from openai import OpenAI
import fitz  # PyMuPDF
from pptx import Presentation
from io import BytesIO, StringIO

# --- Configuration & Secrets ---
# On Streamlit Cloud, these are set in the Sidebar under "Secrets"
GROK_API_KEY = st.secrets.get("GROK_API_KEY", "")
if not GROK_API_KEY:
    st.warning("Grok API Key not found. Please set it in Streamlit Secrets.")

client = OpenAI(
    api_key=GROK_API_KEY,
    base_url="https://api.x.ai/v1",
)

# --- Database Setup ---
def init_db():
    conn = sqlite3.connect("engagement.db")
    c = conn.cursor()
    c.execute('''CREATE TABLE IF NOT EXISTS question_bank 
                 (id INTEGER PRIMARY KEY AUTOINCREMENT, text TEXT, options TEXT, correct_answer TEXT, explanation TEXT)''')
    c.execute('''CREATE TABLE IF NOT EXISTS submissions 
                 (id INTEGER PRIMARY KEY AUTOINCREMENT, name TEXT, email TEXT, score INTEGER, responses TEXT, rating INTEGER, comments TEXT)''')
    conn.commit()
    conn.close()

def db_query(query, params=(), commit=False):
    conn = sqlite3.connect("engagement.db")
    df = None
    try:
        if commit:
            c = conn.cursor()
            c.execute(query, params)
            conn.commit()
        else:
            df = pd.read_sql_query(query, conn, params=params)
    finally:
        conn.close()
    return df

# --- Services ---
def extract_text(uploaded_file):
    text = ""
    file_bytes = uploaded_file.read()
    if uploaded_file.name.endswith(".pdf"):
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        for page in doc:
            text += page.get_text()
    elif uploaded_file.name.endswith(".pptx"):
        prs = Presentation(BytesIO(file_bytes))
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    return text

def generate_questions(slide_text):
    prompt = f"""
    Generate 20 multiple-choice questions based on the following presentation text.
    Each question should have 4 options and 1 correct answer.
    Provide an explanation for the correct answer.
    Format your response as a valid JSON list of objects:
    [
        {{
            "text": "Question text?",
            "options": ["A", "B", "C", "D"],
            "correct_answer": "A",
            "explanation": "Why A is correct."
        }}
    ]
    
    Presentation Text:
    {slide_text}
    """
    
    response = client.chat.completions.create(
        model="grok-beta",
        messages=[
            {"role": "system", "content": "You are a helpful assistant that generates educational quizzes."},
            {"role": "user", "content": prompt}
        ],
        response_format={ "type": "json_object" }
    )
    
    raw_content = response.choices[0].message.content
    # Sometimes models wrap JSON in markdown
    if "```json" in raw_content:
        raw_content = raw_content.split("```json")[1].split("```")[0].strip()
    
    return json.loads(raw_content)

# --- App Layout ---
st.set_page_config(page_title="EngageIQ", page_icon="🚀", layout="wide")
init_db()

# Premium CSS
st.markdown("""
<style>
    .stApp {
        background: radial-gradient(circle at top left, #1e293b, #0f172a);
        color: white;
    }
    .main-header {
        font-size: 3rem;
        font-weight: 800;
        background: linear-gradient(135deg, #38bdf8 0%, #818cf8 100%);
        -webkit-background-clip: text;
        -webkit-text-fill-color: transparent;
        margin-bottom: 2rem;
    }
    .glass-card {
        background: rgba(255, 255, 255, 0.05);
        backdrop-filter: blur(12px);
        border: 1px solid rgba(255, 255, 255, 0.1);
        padding: 2rem;
        border-radius: 1.5rem;
        margin-bottom: 1rem;
    }
</style>
""", unsafe_allow_html=True)

st.markdown('<h1 class="main-header">EngageIQ</h1>', unsafe_allow_html=True)

# Navigation
tabs = st.tabs(["🏠 Home", "🎤 Presenter", "✍️ Join Quiz"])

with tabs[0]:
    st.markdown("""
    ## Elevate Your Presentations
    Transform your slides into interactive learning experiences in seconds.
    
    1. **Upload** your PDF or PPTX slides as a Presenter.
    2. **Grok AI** automatically generates a conceptual question bank.
    3. **Invite** your audience to take a randomized 5-question quiz.
    4. **Analyze** feedback and scores instantly.
    """)

with tabs[1]:
    st.subheader("Presenter Dashboard")
    col1, col2 = st.columns([2, 1])
    
    with col1:
        st.markdown('<div class="glass-card">', unsafe_allow_html=True)
        uploaded_file = st.file_uploader("Upload Presentation Slides", type=["pdf", "pptx"])
        if uploaded_file and st.button("Generate Question Bank"):
            with st.spinner("Grok is analyzing your slides..."):
                try:
                    text = extract_text(uploaded_file)
                    questions_data = generate_questions(text)
                    
                    # Clear old questions optional or append? Let's clear for a fresh session.
                    db_query("DELETE FROM question_bank", commit=True)
                    
                    for q in questions_data:
                        db_query("INSERT INTO question_bank (text, options, correct_answer, explanation) VALUES (?, ?, ?, ?)",
                                 (q["text"], ";".join(q["options"]), q["correct_answer"], q.get("explanation", "")),
                                 commit=True)
                    st.success(f"Successfully generated {len(questions_data)} conceptual questions!")
                except Exception as e:
                    st.error(f"Failed to generate questions: {e}")
        st.markdown('</div>', unsafe_allow_html=True)
    
    with col2:
        st.markdown('<div class="glass-card">', unsafe_allow_html=True)
        st.write("Current Submissions")
        subs_df = db_query("SELECT * FROM submissions")
        if not subs_df.empty:
            st.dataframe(subs_df[["name", "score", "rating"]])
            csv = subs_df.to_csv(index=False).encode('utf-8')
            st.download_button("Download Full Results (CSV)", data=csv, file_name="session_results.csv", mime="text/csv")
        else:
            st.info("No submissions yet.")
        st.markdown('</div>', unsafe_allow_html=True)

with tabs[2]:
    st.subheader("Participant Entrance")
    
    # Check if questions exist
    q_count = db_query("SELECT COUNT(*) as count FROM question_bank")["count"][0]
    
    if q_count == 0:
        st.warning("The presenter hasn't generated any questions yet. Please wait.")
    else:
        if "quiz_started" not in st.session_state:
            st.session_state.quiz_started = False
            st.session_state.answers = {}
            st.session_state.user_data = {"name": "", "email": ""}
            st.session_state.random_questions = []

        if not st.session_state.quiz_started:
            with st.form("user_details"):
                name = st.text_input("Full Name")
                email = st.text_input("Email")
                if st.form_submit_button("Start Knowledge Check"):
                    if name and email:
                        st.session_state.user_data = {"name": name, "email": email}
                        # Select 5 random questions
                        all_qs = db_query("SELECT * FROM question_bank")
                        selected_indices = random.sample(range(len(all_qs)), min(5, len(all_qs)))
                        st.session_state.random_questions = all_qs.iloc[selected_indices].to_dict('records')
                        st.session_state.quiz_started = True
                        st.rerun()
                    else:
                        st.error("Please provide both name and email.")
        else:
            # Quiz UI
            st.write(f"Good luck, {st.session_state.user_data['name']}!")
            
            with st.form("quiz_form"):
                for idx, q in enumerate(st.session_state.random_questions):
                    st.markdown(f"**Q{idx+1}: {q['text']}**")
                    options = q["options"].split(";")
                    st.session_state.answers[q["id"]] = st.radio(f"Select choice for Q{idx+1}", options, key=f"q_{q['id']}", label_visibility="collapsed")
                    st.divider()
                
                # Feedback
                st.subheader("Session Feedback")
                rating = st.select_slider("How would you rate this session?", options=[1, 2, 3, 4, 5], value=5)
                comments = st.text_area("Additional comments (optional)")
                
                if st.form_submit_button("Submit Quiz & Feedback"):
                    score = 0
                    results = []
                    for q in st.session_state.random_questions:
                        user_ans = st.session_state.answers.get(q["id"])
                        is_correct = (user_ans == q["correct_answer"])
                        if is_correct:
                            score += 1
                        results.append({
                            "question": q["text"],
                            "user_answer": user_ans,
                            "correct_answer": q["correct_answer"],
                            "is_correct": is_correct
                        })
                    
                    # Save Submission
                    db_query("INSERT INTO submissions (name, email, score, responses, rating, comments) VALUES (?, ?, ?, ?, ?, ?)",
                             (st.session_state.user_data["name"], st.session_state.user_data["email"], 
                              score, json.dumps(results), rating, comments),
                             commit=True)
                    
                    st.balloons()
                    st.success(f"Done! Your score: {score}/5")
                    st.info("Thank you for your feedback!")
                    
                    # Reset state
                    del st.session_state.quiz_started
                    if st.button("Take Quiz Again"):
                        st.rerun()
